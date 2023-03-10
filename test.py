import collections 
import collections.abc
import fitz

from pptx import Presentation
from pptx.util import Cm

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
print(
    '''
     ____  ____  _____   _____ ___    ____  ____ _____ 
    |  _ \|  _ \|  ___| |_   _/ _ \  |  _ \|  _ \_   _|
    | |_) | | | | |_      | || | | | | |_) | |_) || |  
    |  __/| |_| |  _|     | || |_| | |  __/|  __/ | |  
    |_|   |____/|_|       |_| \___/  |_|   |_|    |_|  
                                                  
--------------------------------------------------------------


    '''

)


pdf_path = input("PDF file name : ")
pdf_doc = fitz.open(pdf_path)

def convert_to_pdf():
    for i in range(pdf_doc.page_count):
        page = pdf_doc.load_page(i)
        pix = page.get_pixmap()
        output = "./output/{}.png".format(i)
        pix.save(output)

convert_to_pdf()

def add_slide_and_image(img_path):
    slide = prs.slides.add_slide(blank_slide_layout) # make new slide
    left = top = Cm(0)
    pic = slide.shapes.add_picture(img_path, left, top, width=Cm(25.4)) # add image
    pic.top = (prs.slide_height - pic.height) // 2

for i in range(pdf_doc.page_count):
    add_slide_and_image("./output/{}.png".format(i))
    

prs.save('output.pptx')